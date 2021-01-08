from pptx import Presentation
import sys
import fitz

try:
    from BeautifulSoup import BeautifulSoup as bs
except ImportError: 
    from bs4 import BeautifulSoup as bs

filename = sys.argv[1]
# prs = Presentation(filename)
# for slide in prs.slides:
#     title = slide.shapes.title
#     print(title)
doc = fitz.open(filename)
for page in doc:
    html = page.getText("html")
    parsed_html = bs(html,"lxml")
    # print(parsed_html)
    print(parsed_html.body.find('p',attrs={'style'}))
    exit(1)
    